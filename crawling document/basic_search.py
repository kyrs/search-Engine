# -*- coding: utf-8 -*-
"""
Created on Wed August 5 14:45:09 2015

@author: shubham
"""
import json 
import requests
from elasticsearch import Elasticsearch
import os
#import slate
from functools import partial
from operator import ne
from pptx import Presentation
import textract
import nltk
from nltk.stem import SnowballStemmer
from nltk.corpus import stopwords
import re
import sys

class Paragraphs:
    def __init__(self, fileobj):
        # self.seq: the underlying line-sequence
        # self.line_num: current index into self.seq (line number)
        # self.para_num: current index into self (paragraph number)
        #import xreadlines
        document = open(fileobj,"r")
        content = document.read()
        content = content.split("\r\n\r\n")
        self.content = filter(partial(ne,""),content)        
        

class crawler (object):
    def __init__(self,book):
        self.document =book
        self.type = book[book.index(".")+1:]
        self.stopwords = stopwords.words('english')
#    def ppt_txt_cnverter(self):
#        prs = Presentation(self.document)
#        text_runs=[]
#        # text_runs will be populated with a list of strings,
#        # one for each text run in presentation
#        for slide in prs.slides:
#            for shape in slide.shapes:
#                if not shape.has_text_frame:
#                    continue
#                for paragraph in shape.text_frame.paragraphs:
#                    #text_runs.append(1)
#                    for run in paragraph.runs:
#                        text_runs.append(run.text)
#        self.doc = text_runs
#
#
#
#        
#    def pdf_txt_converter(self):
#        with open(self.document) as f :
#            self.doc = slate.PDF(f)
#
#
#    def doc_txt_converter(self):
#        (fi, fo, fe) = os.popen3('catdoc -w "%s"' % self.document)
#        fi.close()
#        retval = fo.read()
#        erroroutput = fe.read()
#        fo.close()
#        fe.close()
#        if not erroroutput:
#            self.doc=retval
#        else:
#            raise OSError("Executing the command caused an error: %s" % erroroutput)
#    def paragraph(self):
#        content = self.doc.split("\r\n\r\n")
#        self.doc = filter(partial(ne,""),content)    
        
        
        
    def  crawl(self):
        self.content = textract.process(self.document)
        
        if self.type == "txt":
            try:
                self.data = self.content.lower().split("\r\n\r\n")
                self.data = filter(partial(ne,""),self.data)
            except:
                self.data = self.content.lower().split("\n\n")
                self.data = filter(partial(ne,""),self.data)
        else:
            self. data = self.content.lower().split("\n\n")
            self.data = filter(partial(ne,""),self.data)
    
    
    def json_creater(self):
        self.json = {}
        self.json["Name"] = self.document
        
    def clean_sentence(self,data):
       string = ""
       count = 0
       for i in data :
           if ord(i)<128:
               if ord(i)>65:
                   count+=1
               string+=i
       if count ==0:
           string = ""
       string = re.sub('\\n'," ", string)
       string = re.sub('<(.*?)>'," ",string)
       
       
       return string
       
       
    def data_cleaning(self):
        #self.data.remove(" ")
        #self.data.remove("")
        for i in range(len(self.data)):
            self.data[i] = self.clean_sentence(self.data[i])
        self.data = filter(partial(ne,""), self.data)
        
       
            
                    
                    
            

    def json_converter(self):
        output = {}
        output["Name"]= self.document
        output["type"] = self.type
        paragraph_count = len(self.data)
        output["paragraph"] = paragraph_count
        for i in range(paragraph_count):
            paragraph = "par%d"%(i)
            output[paragraph] = self.data[i]
        self.json = output
        
        


class search(object):
    def __init__(self):
        self.elastic = Elasticsearch([{'host': 'localhost', 'port': 9200}])
        self.response = requests.get('http://localhost:9200')
        
    def check(self):
        print self.response.content
    def crawl(self,**arg):
        url = arg.get('url',"")
        i = int(arg.get('start',2))
        limit = int(arg.get("limit",50))
        while self.response.status_code ==200:
            r =requests.get(url+ str(i))
            self.elastic.index(index="sw",doc_type="people",id =i,
                               body = json.loads(r.content))
            i = i+1
            print i
            if i == limit :
                break

if __name__== "__main__":
    file_name = sys.argv[1]
    document = crawler("./Documents/%s"%(file_name))
    document.crawl()
    document.data_cleaning()
    document.json_converter()
    print document.json